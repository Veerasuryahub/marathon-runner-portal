"""
Marathon Runner Portal — Sprint 1 Review
Professional PowerPoint Generator (20 slides)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
from pptx.enum.dml import MSO_THEME_COLOR
import pptx.oxml.ns as pns
from lxml import etree
import copy

# ── BRAND COLOURS ──────────────────────────────────────────────────────────────
NAVY   = RGBColor(0x0D, 0x2B, 0x55)   # slide background / dark blocks
BLUE   = RGBColor(0x1A, 0x73, 0xE8)   # accent, headers
LBLUE  = RGBColor(0xE8, 0xF0, 0xFE)   # light blue panels
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GREY   = RGBColor(0xF5, 0xF7, 0xFA)
GREEN  = RGBColor(0x34, 0xA8, 0x53)
ORANGE = RGBColor(0xFF, 0x6D, 0x00)
RED    = RGBColor(0xEA, 0x43, 0x35)
LTGREY = RGBColor(0x9E, 0x9E, 0x9E)
ACCENT = RGBColor(0x00, 0xBF, 0xD8)

W = Inches(13.33)   # 16:9 width
H = Inches(7.5)     # 16:9 height

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
blank = prs.slide_layouts[6]   # fully blank layout

# ── HELPERS ────────────────────────────────────────────────────────────────────
def add_rect(slide, x, y, w, h, fill=WHITE, alpha=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    return shape

def add_text(slide, text, x, y, w, h, size=16, bold=False, color=NAVY,
             align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb

def add_title_bar(slide, title, subtitle=""):
    """Navy top bar with white title"""
    bar = add_rect(slide, 0, 0, W, Inches(1.15), fill=NAVY)
    add_text(slide, title, Inches(0.4), Inches(0.08), Inches(10), Inches(0.65),
             size=26, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text(slide, subtitle, Inches(0.4), Inches(0.72), Inches(9), Inches(0.38),
                 size=13, bold=False, color=ACCENT, align=PP_ALIGN.LEFT)

def add_card(slide, x, y, w, h, title, body, title_color=BLUE, bg=LBLUE):
    card = add_rect(slide, x, y, w, h, fill=bg)
    card.line.color.rgb = BLUE
    card.line.width = Pt(0.75)
    add_text(slide, title, x+Inches(0.12), y+Inches(0.08), w-Inches(0.24), Inches(0.3),
             size=12, bold=True, color=title_color, align=PP_ALIGN.LEFT)
    add_text(slide, body,  x+Inches(0.12), y+Inches(0.38), w-Inches(0.24), h-Inches(0.45),
             size=10, bold=False, color=NAVY, align=PP_ALIGN.LEFT)

def add_badge(slide, x, y, text, color=GREEN):
    b = add_rect(slide, x, y, Inches(1.35), Inches(0.28), fill=color)
    b.line.fill.background()
    add_text(slide, text, x+Inches(0.06), y+Inches(0.02), Inches(1.25), Inches(0.26),
             size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

def add_flow_arrow(slide, x, y):
    """Small downward arrow text"""
    add_text(slide, "↓", x, y, Inches(0.3), Inches(0.3), size=14, bold=True, color=BLUE, align=PP_ALIGN.CENTER)

def notes(slide, text):
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = text

def slide_number(slide, num):
    add_text(slide, f"{num} / 20", W-Inches(1.1), H-Inches(0.38), Inches(0.9), Inches(0.3),
             size=9, color=LTGREY, align=PP_ALIGN.RIGHT)

# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 1 — TITLE
# ══════════════════════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(blank)

# Full navy background
add_rect(s1, 0, 0, W, H, fill=NAVY)

# Blue accent stripe
add_rect(s1, 0, Inches(3.9), W, Inches(0.06), fill=BLUE)

# Title
add_text(s1, "MARATHON RUNNER PORTAL",
         Inches(0.7), Inches(1.3), Inches(11.5), Inches(1.1),
         size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Subtitle tag
add_text(s1, "Sprint 1 Review  ·  Java Console Application  ·  Core Java 17",
         Inches(0.7), Inches(2.45), Inches(11.5), Inches(0.5),
         size=17, bold=False, color=ACCENT, align=PP_ALIGN.CENTER)

# Tech pills
pills = [("Core Java 17",""), ("ArrayList / HashMap",""), ("No DB / No Spring",""),
         ("Single Java File",""), ("1505 Lines","")]
px = Inches(1.35)
for label, _ in pills:
    w_pill = Inches(1.9)
    box = add_rect(s1, px, Inches(3.25), w_pill, Inches(0.38), fill=BLUE)
    box.line.fill.background()
    add_text(s1, label, px+Inches(0.06), Inches(3.28), w_pill-Inches(0.1), Inches(0.32),
             size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    px += Inches(2.05)

# Description block
add_text(s1,
    "A role-based Java Console Application for marathon runners and administrators.\n"
    "Runners manage race registrations, training logs, and community forums.\n"
    "Admins control events, users, analytics, and CSV exports.",
    Inches(1.5), Inches(4.15), Inches(10.2), Inches(1.0),
    size=13, color=RGBColor(0xCC,0xDD,0xFF), align=PP_ALIGN.CENTER)

# Developer line
add_text(s1, "TCS ILP  ·  Agile Sprint 1",
         Inches(0.7), Inches(5.8), Inches(11.5), Inches(0.35),
         size=12, color=LTGREY, align=PP_ALIGN.CENTER)

slide_number(s1, 1)
notes(s1, "Introduce the project: Marathon Runner Portal is a Core Java 17 console application. "
      "No database, no Spring Boot, no external libraries. Everything is in one Java file of 1505 lines. "
      "It has two user roles: Runners and Admins. Data lives in ArrayList and HashMap collections.")

# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 2 — PROJECT OVERVIEW
# ══════════════════════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(blank)
add_rect(s2, 0, 0, W, H, fill=GREY)
add_title_bar(s2, "Project Overview", "Problem → Objective → Output")

icons = ["🏃","🎯","💡","📊"]
titles = ["Problem Statement","Objective","Why This Project?","Expected Output"]
bodies = [
    "Runners need a centralized system to register for events, log workouts, and communicate. "
    "Admins need tools to manage races and participants.",
    "Build a role-based Java console app covering user management, race registration, "
    "training tracking, forum, and CSV analytics.",
    "Demonstrates OOP, Java Collections, exception handling, validation, "
    "and layered architecture — all in Core Java with no frameworks.",
    "A fully working console application where runners and admins can log in, "
    "manage data, and export reports to CSV files."
]
for i, (ic, ti, bo) in enumerate(zip(icons, titles, bodies)):
    cx = Inches(0.35 + i * 3.24)
    add_card(s2, cx, Inches(1.35), Inches(3.1), Inches(4.9), f"{ic}  {ti}", bo, bg=WHITE)

slide_number(s2, 2)
notes(s2, "Explain the four key pillars: Problem (need a portal), Objective (build it in Java), "
      "Why (showcase OOP and Java skills), Output (working console app with exports). "
      "Emphasise this is a backend-only app — no UI, no database.")

# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 3 — PROJECT FEATURES
# ══════════════════════════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(blank)
add_rect(s3, 0, 0, W, H, fill=GREY)
add_title_bar(s3, "Project Features", "8 Core Feature Modules")

features = [
    ("👤","User Registration","Email/password/age validation.\nPENDING status until admin verifies."),
    ("🔐","Login & Auth","Role-based dashboard.\n5-attempt lockout for 30 minutes."),
    ("🏁","Race Management","Admin creates/edits/deactivates races.\n7-day future rule, unique names."),
    ("📝","Training Logs","Log workouts with pace calc.\nPace & outlier validations."),
    ("💬","Community Forum","Post, flag spam, URL review.\nAdmin approve / reject."),
    ("📊","Analytics","Dashboard stats for admin & runner.\nRunner training trends."),
    ("📁","CSV Reports","Export users, races, training logs.\nTimestamped CSV files."),
    ("👤","User Profile","View completed races, personal best.\nEdit age & phone only."),
]
cols = 4
for i, (ic, ti, bo) in enumerate(features):
    col = i % cols
    row = i // cols
    cx = Inches(0.3 + col * 3.25)
    cy = Inches(1.35 + row * 2.85)
    add_card(s3, cx, cy, Inches(3.1), Inches(2.65), f"{ic}  {ti}", bo, bg=WHITE)

slide_number(s3, 3)
notes(s3, "Walk through each of the 8 feature modules. Emphasise that all 8 are implemented "
      "as console interactions — no web UI. Each feature maps to a specific class or method in the code.")

# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 4 — AGILE USER STORIES TABLE
# ══════════════════════════════════════════════════════════════════════════════
s4 = prs.slides.add_slide(blank)
add_rect(s4, 0, 0, W, H, fill=GREY)
add_title_bar(s4, "Agile User Stories (US01 – US20)", "Sprint 1 Coverage")

stories = [
    ("US01","User Registration","✅ Implemented","UserMgr.register()"),
    ("US02","Login & Role Dashboard","✅ Implemented","UserMgr.login()"),
    ("US03","Race Listing & Filter","✅ Implemented","viewRaces()"),
    ("US04","Race Registration","✅ Implemented","RegMgr.register()"),
    ("US05","Training Log","✅ Implemented","TrainingMgr.addLog()"),
    ("US06","Training Dashboard","✅ Implemented","viewTrainingDashboard()"),
    ("US07","Community Forum","✅ Implemented","ForumMgr.submit()"),
    ("US08","User Profile","✅ Implemented","viewProfile()"),
    ("US09","Admin Race Management","✅ Implemented","adminRaces()"),
    ("US10","Admin Analytics","✅ Implemented","adminAnalytics()"),
    ("US11","Auth Audit Log","✅ Implemented","logFailedAttempt()"),
    ("US12","Race CRUD Validation","✅ Implemented","RaceMgr.create/edit()"),
    ("US13","Registration Validation","✅ Implemented","RegMgr.register()"),
    ("US14","Training Validation","✅ Implemented","TrainingMgr.addLog()"),
    ("US15","Forum Moderation","✅ Implemented","adminForum()"),
    ("US16","Email Verification","⚠️ Console-Adapted","Admin manual verify"),
    ("US17","Analytics Reports","✅ Implemented","adminAnalytics()"),
    ("US18","Registration Scheduler","⚠️ Console-Adapted","Auto 48h cutoff"),
    ("US19","Role-Based Auth","✅ Implemented","handleLogin() role check"),
    ("US20","Data Export","✅ Implemented","AnalyticsMgr.*"),
]

# Two-column layout
col_w = Inches(6.3)
headers = ["#","User Story","Status","Method"]
col_widths = [Inches(0.55), Inches(2.2), Inches(1.5), Inches(1.9)]

for col_start in range(2):
    ox = Inches(0.25 + col_start * 6.55)
    oy = Inches(1.25)

    # Header row
    hx = ox
    for j, (hdr, cw) in enumerate(zip(headers, col_widths)):
        hb = add_rect(s4, hx, oy, cw, Inches(0.28), fill=NAVY)
        hb.line.fill.background()
        add_text(s4, hdr, hx+Inches(0.05), oy+Inches(0.03), cw-Inches(0.06), Inches(0.23),
                 size=9, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
        hx += cw

    # Rows
    chunk = stories[col_start*10:(col_start+1)*10]
    for ri, row in enumerate(chunk):
        ry = oy + Inches(0.28 + ri*0.27)
        bg = WHITE if ri % 2 == 0 else LBLUE
        rx = ox
        for j, (cell, cw) in enumerate(zip(row, col_widths)):
            cb = add_rect(s4, rx, ry, cw, Inches(0.27), fill=bg)
            cb.line.color.rgb = RGBColor(0xCC,0xCC,0xCC)
            cb.line.width = Pt(0.5)
            clr = GREEN if "✅" in cell else (ORANGE if "⚠️" in cell else NAVY)
            add_text(s4, cell, rx+Inches(0.05), ry+Inches(0.02), cw-Inches(0.06), Inches(0.23),
                     size=8.5, color=clr, align=PP_ALIGN.LEFT)
            rx += cw

slide_number(s4, 4)
notes(s4, "18 of 20 user stories fully implemented. US16 (email verification) adapted to admin manual "
      "verification because SMTP is a web feature. US18 (scheduler) adapted to automatic 48-hour "
      "cutoff calculation using ChronoUnit — no background thread needed in console app.")

# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 5 — SPRINT SCOPE
# ══════════════════════════════════════════════════════════════════════════════
s5 = prs.slides.add_slide(blank)
add_rect(s5, 0, 0, W, H, fill=GREY)
add_title_bar(s5, "Sprint Scope", "What's In · What's Out · Why")

# Left: implemented
add_rect(s5, Inches(0.25), Inches(1.3), Inches(4.1), Inches(5.85), fill=WHITE).line.color.rgb = GREEN
add_text(s5, "✅  IMPLEMENTED", Inches(0.35), Inches(1.38), Inches(3.9), Inches(0.32),
         size=12, bold=True, color=GREEN)
impl = ["User Registration & Login","Role-Based Menus","Race CRUD & Validation",
        "Training Log & Analytics","Community Forum","Admin Approval Workflows",
        "CSV Export Reports","Security Audit Log","In-Memory Collections","Input Validation"]
for i, item in enumerate(impl):
    add_text(s5, f"  •  {item}", Inches(0.35), Inches(1.72+i*0.48), Inches(3.9), Inches(0.42),
             size=10.5, color=NAVY)

# Middle: console-adapted
add_rect(s5, Inches(4.6), Inches(1.3), Inches(3.6), Inches(5.85), fill=WHITE).line.color.rgb = ORANGE
add_text(s5, "⚠️  CONSOLE-ADAPTED", Inches(4.7), Inches(1.38), Inches(3.4), Inches(0.32),
         size=12, bold=True, color=ORANGE)
adapted = [
    ("Email Verification","Admin manual verify"),
    ("Background Scheduler","Auto 48h time check"),
]
for i, (feat, how) in enumerate(adapted):
    add_text(s5, f"  ✦  {feat}", Inches(4.7), Inches(1.75+i*0.95), Inches(3.4), Inches(0.28),
             size=11, bold=True, color=NAVY)
    add_text(s5, f"     ↳ {how}", Inches(4.7), Inches(2.05+i*0.95), Inches(3.4), Inches(0.28),
             size=10, italic=True, color=LTGREY)

# Right: not implemented
add_rect(s5, Inches(8.45), Inches(1.3), Inches(4.6), Inches(5.85), fill=WHITE).line.color.rgb = RED
add_text(s5, "🚫  NOT IN SCOPE (Web-Only)", Inches(8.55), Inches(1.38), Inches(4.4), Inches(0.32),
         size=12, bold=True, color=RED)
nscope = [
    ("MySQL / MongoDB","Needs JDBC + DB server"),
    ("REST APIs","Needs Servlet / Spring"),
    ("JWT Auth","Needs HTTP headers"),
    ("SMTP Email","Needs mail server"),
    ("Browser Sessions","Needs cookies/web"),
    ("Background Threads","Needs scheduler daemon"),
    ("HTML / CSS","Needs browser"),
    ("PDF Export","Needs Apache POI"),
]
for i, (feat, reason) in enumerate(nscope):
    add_text(s5, f"  ✦  {feat}", Inches(8.55), Inches(1.75+i*0.58), Inches(4.3), Inches(0.26),
             size=10, bold=True, color=NAVY)
    add_text(s5, f"     {reason}", Inches(8.55), Inches(2.02+i*0.58), Inches(4.3), Inches(0.24),
             size=8.5, italic=True, color=LTGREY)

slide_number(s5, 5)
notes(s5, "Explain the three zones: Green = fully done, Orange = done in a console-appropriate way, "
      "Red = web-only features intentionally excluded. This shows the team understands scope boundaries.")

# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 6 — SYSTEM ARCHITECTURE
# ══════════════════════════════════════════════════════════════════════════════
s6 = prs.slides.add_slide(blank)
add_rect(s6, 0, 0, W, H, fill=GREY)
add_title_bar(s6, "Complete System Architecture", "Layered Architecture — Console MVC Pattern")

# Draw architecture layers
layers = [
    (NAVY, WHITE,  "PRESENTATION LAYER",      "Console Menus  ·  printBanner()  ·  showMainMenu()  ·  showRunnerMenu()  ·  showAdminMenu()"),
    (BLUE, WHITE,  "AUTHENTICATION LAYER",     "UserMgr.login()  ·  isLocked()  ·  Role Check  ·  currentUser session"),
    (RGBColor(0x00,0x89,0x7B), WHITE, "BUSINESS LOGIC LAYER",
     "UserMgr  ·  RaceMgr  ·  RegMgr  ·  TrainingMgr  ·  ForumMgr  ·  AnalyticsMgr"),
    (RGBColor(0x55,0x57,0xBF), WHITE, "VALIDATION LAYER",
     "Email · Password · Age · Phone · Date · Pace · Distance · Duplicate · Capacity"),
    (RGBColor(0x8E,0x24,0xAA), WHITE, "DATA / STORAGE LAYER",
     "ArrayList<User>  ·  ArrayList<Race>  ·  ArrayList<TrainingLog>  ·  HashMap<String,User>  ·  HashSet<String>"),
    (LTGREY, NAVY, "ENTITIES LAYER",
     "User  ·  Race  ·  TrainingLog  ·  ForumPost  ·  Registration"),
]

for i, (bg, fg, title, detail) in enumerate(layers):
    ly = Inches(1.3 + i * 0.97)
    block = add_rect(s6, Inches(0.5), ly, Inches(12.3), Inches(0.85), fill=bg)
    block.line.fill.background()
    add_text(s6, title, Inches(0.65), ly+Inches(0.06), Inches(3.5), Inches(0.35),
             size=11, bold=True, color=fg)
    add_text(s6, detail, Inches(4.3), ly+Inches(0.06), Inches(8.3), Inches(0.75),
             size=10, color=fg)
    if i < len(layers)-1:
        add_text(s6, "↕", Inches(6.5), Inches(1.3 + i*0.97 + 0.87), Inches(0.4), Inches(0.1),
                 size=11, bold=True, color=BLUE, align=PP_ALIGN.CENTER)

slide_number(s6, 6)
notes(s6, "Six-layer architecture from top to bottom: Presentation (menus), Authentication (login/roles), "
      "Business Logic (6 manager classes), Validation (all rules), Data Storage (collections), "
      "Entities (5 data classes). Each layer talks only to the layer below it.")

# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 7 — PROGRAM FLOW CHART
# ══════════════════════════════════════════════════════════════════════════════
s7 = prs.slides.add_slide(blank)
add_rect(s7, 0, 0, W, H, fill=GREY)
add_title_bar(s7, "Complete Program Flow", "From startup to output — step by step")

def flow_box(slide, x, y, w, h, text, bg=BLUE, fg=WHITE, sz=10.5, bold=True):
    b = add_rect(slide, x, y, w, h, fill=bg)
    b.line.color.rgb = NAVY
    b.line.width = Pt(0.75)
    add_text(slide, text, x+Inches(0.06), y+Inches(0.05), w-Inches(0.1), h-Inches(0.08),
             size=sz, bold=bold, color=fg, align=PP_ALIGN.CENTER)

def flow_diamond(slide, x, y, w, h, text):
    # Simulate diamond with rotated rectangle
    b = add_rect(slide, x, y, w, h, fill=ORANGE)
    b.line.color.rgb = NAVY
    b.line.width = Pt(0.75)
    add_text(slide, text, x+Inches(0.05), y+Inches(0.05), w-Inches(0.08), h-Inches(0.08),
             size=9.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

cx = Inches(0.3)
# Main vertical flow (left column)
flow_box(s7, cx, Inches(1.3), Inches(2.5), Inches(0.45), "▶  JVM → main()", bg=NAVY)
add_text(s7, "↓", cx+Inches(1.0), Inches(1.77), Inches(0.5), Inches(0.2), size=14, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
flow_box(s7, cx, Inches(1.97), Inches(2.5), Inches(0.45), "loadSeedData()", bg=RGBColor(0x00,0x89,0x7B))
add_text(s7, "↓", cx+Inches(1.0), Inches(2.44), Inches(0.5), Inches(0.2), size=14, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
flow_box(s7, cx, Inches(2.64), Inches(2.5), Inches(0.45), "printBanner()")
add_text(s7, "↓", cx+Inches(1.0), Inches(3.11), Inches(0.5), Inches(0.2), size=14, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
flow_box(s7, cx, Inches(3.3), Inches(2.5), Inches(0.45), "showMainMenu()")
add_text(s7, "↓", cx+Inches(1.0), Inches(3.77), Inches(0.5), Inches(0.2), size=14, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
flow_diamond(s7, cx, Inches(3.97), Inches(2.5), Inches(0.55), "User Choice?")

# Branches
add_text(s7, "[1] Register →", Inches(2.95), Inches(4.04), Inches(1.3), Inches(0.3), size=9, color=GREEN, bold=True)
add_text(s7, "[2] Login →", Inches(2.95), Inches(4.34), Inches(1.3), Inches(0.3), size=9, color=BLUE, bold=True)
add_text(s7, "[0] Exit →", Inches(2.95), Inches(4.64), Inches(1.3), Inches(0.3), size=9, color=RED, bold=True)

# Register branch
flow_box(s7, Inches(4.4), Inches(3.97), Inches(2.0), Inches(0.38), "handleRegister()", bg=GREEN)
add_text(s7, "→", Inches(6.45), Inches(4.04), Inches(0.3), Inches(0.3), size=14, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
flow_box(s7, Inches(6.8), Inches(3.97), Inches(2.0), Inches(0.38), "UserMgr.register()", bg=GREEN)
add_text(s7, "→", Inches(8.85), Inches(4.04), Inches(0.3), Inches(0.3), size=14, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
flow_box(s7, Inches(9.2), Inches(3.97), Inches(2.0), Inches(0.38), "Validate → Save to List", bg=GREEN)
add_text(s7, "→", Inches(11.25), Inches(4.04), Inches(0.3), Inches(0.3), size=14, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
flow_box(s7, Inches(11.6), Inches(3.97), Inches(1.5), Inches(0.38), "PENDING Status", bg=GREEN)

# Login branch
flow_box(s7, Inches(4.4), Inches(4.43), Inches(2.0), Inches(0.38), "handleLogin()", bg=BLUE)
add_text(s7, "→", Inches(6.45), Inches(4.5), Inches(0.3), Inches(0.3), size=14, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
flow_diamond(s7, Inches(6.8), Inches(4.43), Inches(1.8), Inches(0.55), "Role?")

# Admin / Runner sub-branches
add_text(s7, "ADMIN ↘", Inches(8.7), Inches(4.43), Inches(1.1), Inches(0.28), size=9, bold=True, color=NAVY)
add_text(s7, "RUNNER ↗", Inches(8.7), Inches(4.71), Inches(1.1), Inches(0.28), size=9, bold=True, color=BLUE)
flow_box(s7, Inches(9.9), Inches(4.25), Inches(1.6), Inches(0.38), "showAdminMenu()", bg=NAVY)
flow_box(s7, Inches(9.9), Inches(4.75), Inches(1.6), Inches(0.38), "showRunnerMenu()", bg=BLUE)
add_text(s7, "→ 7 options", Inches(11.55), Inches(4.32), Inches(1.0), Inches(0.28), size=8.5, color=NAVY, bold=True)
add_text(s7, "→ 10 options", Inches(11.55), Inches(4.82), Inches(1.0), Inches(0.28), size=8.5, color=BLUE, bold=True)

# Bottom flow
flow_box(s7, Inches(4.4), Inches(5.5), Inches(2.0), Inches(0.45), "Manager Class\n(Validate)", bg=RGBColor(0x55,0x57,0xBF))
add_text(s7, "→", Inches(6.45), Inches(5.62), Inches(0.4), Inches(0.3), size=14, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
flow_box(s7, Inches(6.9), Inches(5.5), Inches(2.0), Inches(0.45), "Collection\n(Store/Update)", bg=RGBColor(0x8E,0x24,0xAA))
add_text(s7, "→", Inches(8.95), Inches(5.62), Inches(0.4), Inches(0.3), size=14, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
flow_box(s7, Inches(9.4), Inches(5.5), Inches(2.0), Inches(0.45), "Console Output\nok() / err() / warn()", bg=RGBColor(0x34,0xA8,0x53))

slide_number(s7, 7)
notes(s7, "Walk through the flow: JVM starts main(), seed data is loaded, banner prints, "
      "then the infinite while loop keeps showing the main menu. Register goes to UserMgr, "
      "Login checks role and routes to Admin or Runner menu. Every user action calls a manager "
      "which validates, stores to a collection, then displays output.")

# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 8 — CODE STRUCTURE
# ══════════════════════════════════════════════════════════════════════════════
s8 = prs.slides.add_slide(blank)
add_rect(s8, 0, 0, W, H, fill=GREY)
add_title_bar(s8, "Complete Code Structure", "MarathonPortal.java — 7 Sections · 1505 Lines")

sections = [
    ("Section 1","Lines 26–191","Data Entities",
     "User  ·  Race  ·  TrainingLog  ·  ForumPost  ·  Registration",
     NAVY),
    ("Section 2","Lines 193–223","In-Memory Database",
     "ArrayList · HashMap · HashSet · ID Counters · Scanner · Formatters",
     BLUE),
    ("Section 3","Lines 225–556","Business Logic Managers",
     "UserMgr · RaceMgr · RegMgr · TrainingMgr · ForumMgr · AnalyticsMgr",
     RGBColor(0x00,0x89,0x7B)),
    ("Section 4","Lines 606–754","Main Menus",
     "showMainMenu() · handleRegister() · handleLogin() · showRunnerMenu() · showAdminMenu()",
     RGBColor(0x55,0x57,0xBF)),
    ("Section 5","Lines 756–1104","Runner Features (10 options)",
     "viewProfile · viewRaces · registerRace · addTraining · viewDashboard · editTraining · viewForum · postForum · exportTraining",
     RGBColor(0x8E,0x24,0xAA)),
    ("Section 6","Lines 1106–1424","Admin Features (7 options)",
     "adminUsers · adminRaces · adminForum · adminAnalytics · adminExports · adminSecurityLog · adminMarkCompleted",
     ORANGE),
    ("Section 7","Lines 1426–1503","Helper IO Functions",
     "readStr · readInt · readDouble · ok · err · warn · pause · line · printBanner · findById",
     RGBColor(0x34,0xA8,0x53)),
]

for i, (sec, lines, title, detail, bg) in enumerate(sections):
    y = Inches(1.3 + i * 0.845)
    block = add_rect(s8, Inches(0.25), y, Inches(1.35), Inches(0.77), fill=bg)
    block.line.fill.background()
    add_text(s8, sec, Inches(0.28), y+Inches(0.04), Inches(1.3), Inches(0.3),
             size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s8, lines, Inches(0.28), y+Inches(0.37), Inches(1.3), Inches(0.28),
             size=8, color=RGBColor(0xCC,0xDD,0xFF), align=PP_ALIGN.CENTER)

    add_rect(s8, Inches(1.68), y, Inches(2.5), Inches(0.77), fill=bg)
    add_text(s8, title, Inches(1.72), y+Inches(0.2), Inches(2.4), Inches(0.4),
             size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    info = add_rect(s8, Inches(4.25), y, Inches(8.8), Inches(0.77), fill=LBLUE)
    info.line.color.rgb = bg
    info.line.width = Pt(1.0)
    add_text(s8, detail, Inches(4.35), y+Inches(0.18), Inches(8.6), Inches(0.45),
             size=9.5, color=NAVY, align=PP_ALIGN.LEFT)

slide_number(s8, 8)
notes(s8, "MarathonPortal.java is divided into 7 clearly labelled sections using comment banners. "
      "Section 1 = entities (data classes). Section 2 = in-memory DB. Section 3 = all business rules. "
      "Sections 4-6 = all menus. Section 7 = input/output helpers. Each section has a clear responsibility.")

# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 9 — METHOD CALL HIERARCHY
# ══════════════════════════════════════════════════════════════════════════════
s9 = prs.slides.add_slide(blank)
add_rect(s9, 0, 0, W, H, fill=GREY)
add_title_bar(s9, "Method Call Hierarchy", "Complete method tree — every function shown")

# Tree diagram using text boxes
tree_data = [
    (0, "main()", NAVY),
    (1, "loadSeedData()  ·  printBanner()  ·  showMainMenu()", BLUE),
    (2, "handleRegister()  →  UserMgr.register()  →  checkPassword()", GREEN),
    (2, "handleLogin()  →  UserMgr.login()  →  isLocked()  →  logFailedAttempt()", BLUE),
    (2, "System.exit(0)", RED),
    (1, "showRunnerMenu()", BLUE),
    (2, "viewProfile()  →  findRaceById()  ·  Personal Bests Map", RGBColor(0x55,0x57,0xBF)),
    (2, "viewRaces()  →  Race filters  ·  signedUp HashSet", RGBColor(0x55,0x57,0xBF)),
    (2, "registerRace()  →  RegMgr.register()  →  logsCount  ·  latePenalty", RGBColor(0x55,0x57,0xBF)),
    (2, "addTraining()  →  TrainingMgr.addLog()  →  pace calc  ·  outlier flag", RGBColor(0x55,0x57,0xBF)),
    (2, "viewTrainingDashboard()  →  stats calc  ·  sort lambdas", RGBColor(0x55,0x57,0xBF)),
    (2, "editTraining()  →  TrainingMgr.editLog()  →  6-month lock", RGBColor(0x55,0x57,0xBF)),
    (2, "postForum()  →  ForumMgr.submit()  →  BLACKLIST scan  ·  link detect", RGBColor(0x55,0x57,0xBF)),
    (2, "exportTraining()  →  AnalyticsMgr.exportTraining()  →  CSV write", RGBColor(0x55,0x57,0xBF)),
    (1, "showAdminMenu()", NAVY),
    (2, "adminUsers()  →  list all  ·  verify PENDING runner", ORANGE),
    (2, "adminRaces()  →  RaceMgr.create/edit/deactivate()", ORANGE),
    (2, "adminForum()  →  approve/reject  ·  ForumMgr.cleanup()", ORANGE),
    (2, "adminAnalytics()  →  stream().filter().count()", ORANGE),
    (2, "adminExports()  →  AnalyticsMgr.exportUsers/Races/Logs()", ORANGE),
    (2, "adminSecurityLog()  →  BufferedReader  ·  security.log", ORANGE),
    (2, "adminMarkCompleted()  →  reg.status='COMPLETED'", ORANGE),
]

indent_w = [Inches(0), Inches(0.35), Inches(0.7)]
prefix   = ["", "├── ", "│   └── "]
sz_map   = [12, 10.5, 9.5]

for i, (lvl, text, color) in enumerate(tree_data):
    y = Inches(1.27 + i * 0.285)
    add_text(s9, prefix[lvl] + text,
             Inches(0.35) + indent_w[lvl], y,
             Inches(12.5) - indent_w[lvl], Inches(0.27),
             size=sz_map[lvl], bold=(lvl==0), color=color)

slide_number(s9, 9)
notes(s9, "This tree shows every method and which methods call which others. "
      "main() is the root. It branches into main menu → register, login, exit. "
      "Login branches to Runner Menu (10 methods) and Admin Menu (7 methods). "
      "Each leaf method calls a Manager class which handles validation and storage.")

# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 10 — MODULE RELATIONSHIP
# ══════════════════════════════════════════════════════════════════════════════
s10 = prs.slides.add_slide(blank)
add_rect(s10, 0, 0, W, H, fill=GREY)
add_title_bar(s10, "Module Relationship Diagram", "How the 6 Manager Modules connect")

# Central node
cx_c, cy_c = Inches(6.5), Inches(4.1)
centre = add_rect(s10, cx_c-Inches(1.05), cy_c-Inches(0.4), Inches(2.1), Inches(0.8), fill=NAVY)
centre.line.fill.background()
add_text(s10, "MarathonPortal\nmain()", cx_c-Inches(1.05)+Inches(0.05), cy_c-Inches(0.4)+Inches(0.04),
         Inches(2.0), Inches(0.72), size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

modules = [
    (Inches(1.1), Inches(1.5), "UserMgr", "register()\nlogin()\ncheckPassword()\nlogFailedAttempt()", BLUE),
    (Inches(5.3), Inches(1.3), "RaceMgr", "create()\nedit()\ndeactivate()", GREEN),
    (Inches(9.8), Inches(1.5), "RegMgr", "register()\nvalidateEligibility()\nlatePenalty()", ORANGE),
    (Inches(1.1), Inches(5.6), "TrainingMgr", "addLog()\neditLog()\npace calc\noutlier flag", RGBColor(0x8E,0x24,0xAA)),
    (Inches(5.3), Inches(5.8), "ForumMgr", "submit()\ncleanup()\nblacklist scan\nURL detect", RED),
    (Inches(9.8), Inches(5.6), "AnalyticsMgr", "exportTraining()\nexportUsers()\nexportRaces()\nexportAllLogs()", RGBColor(0x00,0x89,0x7B)),
]

for mx, my, title, methods, color in modules:
    mod = add_rect(s10, mx, my, Inches(2.8), Inches(1.35), fill=color)
    mod.line.fill.background()
    add_text(s10, title, mx+Inches(0.08), my+Inches(0.05), Inches(2.65), Inches(0.32),
             size=12, bold=True, color=WHITE)
    add_text(s10, methods, mx+Inches(0.08), my+Inches(0.38), Inches(2.65), Inches(0.9),
             size=9, color=WHITE)

# Arrow descriptions between modules
adds = [
    (Inches(3.2), Inches(3.0), "↗ User Auth"),
    (Inches(6.5), Inches(2.6), "↑ Race Events"),
    (Inches(9.8), Inches(3.1), "↖ Signup"),
    (Inches(3.2), Inches(5.3), "↘ Training"),
    (Inches(6.5), Inches(5.5), "↓ Forum Posts"),
    (Inches(9.9), Inches(5.1), "↙ Reports"),
]
for ax, ay, albl in adds:
    add_text(s10, albl, ax, ay, Inches(1.5), Inches(0.3), size=9, bold=True, color=LTGREY)

# Shared collections note
add_rect(s10, Inches(0.3), Inches(7.0), Inches(12.6), Inches(0.32), fill=LBLUE).line.fill.background()
add_text(s10, "All modules share:  users (ArrayList)  ·  races (ArrayList)  ·  trainingLogs (ArrayList)  ·  "
              "registrations (ArrayList)  ·  forumPosts (ArrayList)  ·  userByEmail (HashMap)  ·  usedEmails (HashSet)",
         Inches(0.35), Inches(7.04), Inches(12.5), Inches(0.28),
         size=9, color=NAVY, align=PP_ALIGN.CENTER)

slide_number(s10, 10)
notes(s10, "Six manager modules surround the central main() coordinator. "
      "All modules share the same static collections — no module has its own private storage. "
      "UserMgr handles auth, RaceMgr handles events, RegMgr handles signups, "
      "TrainingMgr handles workouts, ForumMgr handles posts, AnalyticsMgr handles exports.")

# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 11 — USER STORY MAPPING TABLE
# ══════════════════════════════════════════════════════════════════════════════
s11 = prs.slides.add_slide(blank)
add_rect(s11, 0, 0, W, H, fill=GREY)
add_title_bar(s11, "User Story → Code Mapping", "Each story traced to method + validation + collection")

mapping = [
    ("US01","UserMgr.register()","email/pwd/age/phone/dup check","users (ArrayList), usedEmails (HashSet), userByEmail (HashMap)"),
    ("US02","UserMgr.login()","lockout, pending, wrong pwd","userByEmail (HashMap)"),
    ("US03","viewRaces()","active, 48h open, not signed up","races (ArrayList), registrations (ArrayList)"),
    ("US04","RegMgr.register()","≥5 logs, no double-book, capacity","registrations (ArrayList), trainingLogs (ArrayList)"),
    ("US05","TrainingMgr.addLog()","no future, 1-100km, pace 3-10","trainingLogs (ArrayList)"),
    ("US06","viewTrainingDashboard()","avg/best/pace stats + sort","trainingLogs (filtered list)"),
    ("US07","ForumMgr.submit()","banned words, link detect, ban count","forumPosts (ArrayList)"),
    ("US08","viewProfile()","COMPLETED regs, personal best","registrations (ArrayList), HashMap<dist,time>"),
    ("US09","adminRaces()  RaceMgr.*","7-day future, 24h edit lock, unique name","races (ArrayList)"),
    ("US10","adminAnalytics()","stream().filter().count()","All collections (read-only)"),
    ("US11","logFailedAttempt()","every wrong pwd → security.log","File I/O, user.failedAttempts"),
    ("US12","RaceMgr.create/edit()","7-day future, unique per year, 24h lock","races (ArrayList)"),
    ("US13","RegMgr.register()","≥5 logs, full capacity, late penalty","registrations, trainingLogs"),
    ("US14","TrainingMgr.addLog()","pace 3-10, distance 1-100, 10x flag","trainingLogs (ArrayList)"),
    ("US15","adminForum()  ForumMgr.cleanup()","approve/reject/1-year cleanup","forumPosts (ArrayList)"),
    ("US17","adminAnalytics()  viewTrainingDashboard()","runner stats, injury flag","trainingLogs, users"),
    ("US19","handleLogin() role check","ADMIN vs RUNNER routing","currentUser.role"),
    ("US20","AnalyticsMgr.*  adminExports()","export all data types","All collections → CSV files"),
]

# Header
headers = ["US #", "Method(s)", "Key Validation", "Collection Used"]
hwidths = [Inches(0.7), Inches(3.2), Inches(3.8), Inches(5.3)]
hx = Inches(0.25)
for h, hw in zip(headers, hwidths):
    hb = add_rect(s11, hx, Inches(1.27), hw, Inches(0.28), fill=NAVY)
    hb.line.fill.background()
    add_text(s11, h, hx+Inches(0.05), Inches(1.29), hw-Inches(0.08), Inches(0.24),
             size=9, bold=True, color=WHITE)
    hx += hw

for ri, (us, meth, val, col) in enumerate(mapping):
    ry = Inches(1.55 + ri * 0.302)
    bg = WHITE if ri % 2 == 0 else LBLUE
    rx = Inches(0.25)
    row_vals = [us, meth, val, col]
    for cell, hw in zip(row_vals, hwidths):
        cb = add_rect(s11, rx, ry, hw, Inches(0.295), fill=bg)
        cb.line.color.rgb = RGBColor(0xCC,0xCC,0xCC)
        cb.line.width = Pt(0.4)
        clr = BLUE if ri == 0 and cell.startswith("US") else NAVY
        add_text(s11, cell, rx+Inches(0.04), ry+Inches(0.03), hw-Inches(0.06), Inches(0.235),
                 size=8.0, color=NAVY, align=PP_ALIGN.LEFT)
        rx += hw

slide_number(s11, 11)
notes(s11, "Every user story traced to the exact Java method, the validations applied, "
      "and which collection stores the data. This demonstrates complete traceability from requirement to code.")

# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 12 — VALIDATION FLOW
# ══════════════════════════════════════════════════════════════════════════════
s12 = prs.slides.add_slide(blank)
add_rect(s12, 0, 0, W, H, fill=GREY)
add_title_bar(s12, "Validation Flow", "Every validation rule — how errors are caught and reported")

# Three validation trees
val_groups = [
    ("👤  REGISTRATION VALIDATION", GREEN, [
        "Username ≥3 chars  [a-zA-Z0-9_]+",
        "Email: @gmail.com or @yahoo.com only",
        "Email: must not exist in usedEmails HashSet",
        "Password: ≥12 chars + UPPER + DIGIT + SPECIAL",
        "Age: between 18 and 40",
        "Phone: exactly 10 digits  \\d{10}",
        "→  throw Exception(message) if any fails",
        "→  ok() on success + PENDING status",
    ]),
    ("🏁  RACE & REGISTRATION VALIDATION", ORANGE, [
        "Race date: ≥7 days in future (ChronoUnit.DAYS)",
        "Race name: unique per year (loop + equalsIgnoreCase)",
        "Race capacity: ≥10 participants",
        "Race edit: blocked if <24h to start",
        "Registration: race must be ACTIVE",
        "Registration: must be open (≥48h, isRegistrationOpen())",
        "Registration: race not full (isFull())",
        "Registration: no double-booking same date",
        "Registration: ≥5 training logs in 30 days",
        "Late penalty: −50 pts if within 7 days",
    ]),
    ("📝  TRAINING & FORUM VALIDATION", RGBColor(0x8E,0x24,0xAA), [
        "Training: no future dates",
        "Training: distance 1.0 – 100.0 km",
        "Training: pace 3.0 – 10.0 min/km",
        "Training: no duplicate (same date + distance)",
        "Training: edit blocked if log >6 months old",
        "Training: auto-flag if distance ≥10× average",
        "Forum: not banned (forumBanned == false)",
        "Forum: title ≥5 chars, content ≥10 chars",
        "Forum: BLACKLIST word scan → FLAGGED",
        "Forum: 3 flagged posts → forumBanned = true",
    ]),
]

col_w = Inches(4.2)
for gi, (title, color, rules) in enumerate(val_groups):
    gx = Inches(0.25 + gi * 4.35)
    header = add_rect(s12, gx, Inches(1.3), col_w, Inches(0.38), fill=color)
    header.line.fill.background()
    add_text(s12, title, gx+Inches(0.08), Inches(1.33), col_w-Inches(0.1), Inches(0.32),
             size=10.5, bold=True, color=WHITE)
    for ri, rule in enumerate(rules):
        ry = Inches(1.72 + ri * 0.53)
        bg = WHITE if ri % 2 == 0 else LBLUE
        rb = add_rect(s12, gx, ry, col_w, Inches(0.5), fill=bg)
        rb.line.color.rgb = color
        rb.line.width = Pt(0.5)
        add_text(s12, rule, gx+Inches(0.1), ry+Inches(0.08), col_w-Inches(0.15), Inches(0.38),
                 size=9, color=NAVY)

slide_number(s12, 12)
notes(s12, "Three groups of validations: Registration (6 rules), Race/Registration (10 rules), "
      "Training/Forum (10 rules). Every rule throws a typed Exception with a clear error message. "
      "The menu method catches it and displays with err(). User is prompted to retry.")

# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 13 — COLLECTIONS
# ══════════════════════════════════════════════════════════════════════════════
s13 = prs.slides.add_slide(blank)
add_rect(s13, 0, 0, W, H, fill=GREY)
add_title_bar(s13, "Java Collections Used", "Why each collection was chosen")

collections = [
    ("ArrayList<User>", NAVY, "users",
     "• Ordered list of all accounts\n• Iterable for admin user list\n• Add: O(1) amortised\n• Search: O(n) loop"),
    ("ArrayList<Race>", BLUE, "races",
     "• All race events in order\n• Filterable for runner view\n• Add: O(1)\n• Supports removeIf()"),
    ("ArrayList<TrainingLog>", GREEN, "trainingLogs",
     "• All workouts for all runners\n• Filter by runnerId\n• Sortable with lambda\n• Add: O(1)"),
    ("ArrayList<Registration>", ORANGE, "registrations",
     "• All race signups\n• Double-booking check\n• Eligibility count source\n• Completion tracking"),
    ("ArrayList<ForumPost>", RGBColor(0x8E,0x24,0xAA), "forumPosts",
     "• All forum posts\n• Filter by status\n• removeIf() for cleanup\n• Admin review queue"),
    ("HashMap<String,User>", RED, "userByEmail",
     "• Key: email → Value: User\n• O(1) average lookup\n• Used every login call\n• Prevents O(n) scan"),
    ("HashSet<String>", RGBColor(0x00,0x89,0x7B), "usedEmails",
     "• Set of registered emails\n• O(1) duplicate check\n• contains() on every signup\n• Prevents duplicate accounts"),
]

cols = 4
for i, (name, color, var, desc) in enumerate(collections):
    col = i % cols
    row = i // cols
    cx = Inches(0.25 + col * 3.25)
    cy = Inches(1.35 + row * 2.75)
    cw, ch = Inches(3.0), Inches(2.55)
    header = add_rect(s13, cx, cy, cw, Inches(0.45), fill=color)
    header.line.fill.background()
    add_text(s13, name, cx+Inches(0.08), cy+Inches(0.06), cw-Inches(0.12), Inches(0.32),
             size=10, bold=True, color=WHITE)
    body = add_rect(s13, cx, cy+Inches(0.45), cw, ch-Inches(0.45), fill=WHITE)
    body.line.color.rgb = color
    body.line.width = Pt(0.75)
    add_text(s13, f"var: {var}", cx+Inches(0.08), cy+Inches(0.48), cw-Inches(0.12), Inches(0.25),
             size=8.5, bold=True, italic=True, color=LTGREY)
    add_text(s13, desc, cx+Inches(0.08), cy+Inches(0.72), cw-Inches(0.12), ch-Inches(0.8),
             size=9.5, color=NAVY)

slide_number(s13, 13)
notes(s13, "Seven collections: five ArrayLists (users, races, training, registrations, forum posts), "
      "one HashMap (email→user for O(1) login lookup), one HashSet (email set for O(1) duplicate check). "
      "ArrayList chosen for ordered, iterable storage. HashMap/HashSet chosen for constant-time operations.")

# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 14 — OOP CONCEPTS
# ══════════════════════════════════════════════════════════════════════════════
s14 = prs.slides.add_slide(blank)
add_rect(s14, 0, 0, W, H, fill=GREY)
add_title_bar(s14, "OOP Concepts Demonstrated", "Where every concept appears in the code")

oops = [
    ("🔒  Encapsulation", NAVY,
     "User, Race, TrainingLog group related fields.\n"
     "UserMgr encapsulates all auth logic.\n"
     "Menus call managers without knowing internals."),
    ("🎭  Abstraction", BLUE,
     "showRunnerMenu() calls addTraining()\n"
     "which calls TrainingMgr.addLog().\n"
     "Each level hides its complexity."),
    ("🔄  Polymorphism", GREEN,
     "Lambda comparators passed to sort():\n"
     "(a,b)->b.date.compareTo(a.date)\n"
     "Same sort() method, different behavior."),
    ("📦  Static Members", ORANGE,
     "Collections & ID counters are static.\n"
     "Shared across all methods in the class.\n"
     "Acts as the application's data store."),
    ("🏗️  Inner Classes", RGBColor(0x8E,0x24,0xAA),
     "UserMgr, RaceMgr etc. are static inner\n"
     "classes of MarathonPortal.\n"
     "Keeps code in one compilable file."),
    ("⚠️  Exception Handling", RED,
     "All validations throw Exception.\n"
     "Menus use try-catch to display err().\n"
     "Specific messages guide users."),
    ("📁  File Handling", RGBColor(0x00,0x89,0x7B),
     "try-with-resources PrintWriter (CSV export).\n"
     "BufferedReader reads security.log.\n"
     "File.mkdirs() creates output folders."),
    ("⏰  Java Time API", RGBColor(0x34,0x69,0xEB),
     "LocalDate, LocalDateTime, ChronoUnit.\n"
     "Race window checks, lockout timers,\n"
     "training age calculations."),
]

cols = 4
for i, (title, color, desc) in enumerate(oops):
    col = i % cols
    row = i // cols
    cx = Inches(0.25 + col * 3.25)
    cy = Inches(1.35 + row * 2.75)
    cw, ch = Inches(3.0), Inches(2.55)
    header = add_rect(s14, cx, cy, cw, Inches(0.42), fill=color)
    header.line.fill.background()
    add_text(s14, title, cx+Inches(0.08), cy+Inches(0.06), cw-Inches(0.12), Inches(0.32),
             size=10, bold=True, color=WHITE)
    body = add_rect(s14, cx, cy+Inches(0.42), cw, ch-Inches(0.42), fill=WHITE)
    body.line.color.rgb = color
    body.line.width = Pt(0.75)
    add_text(s14, desc, cx+Inches(0.08), cy+Inches(0.5), cw-Inches(0.12), ch-Inches(0.55),
             size=10, color=NAVY)

slide_number(s14, 14)
notes(s14, "8 OOP/Java concepts: Encapsulation (classes group data+logic), Abstraction (layered calls), "
      "Polymorphism (lambda comparators), Static members (shared collections), Inner classes (single-file design), "
      "Exception Handling (all validation throws), File Handling (CSV+log), Java Time API (all date logic).")

# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 15 — TESTING
# ══════════════════════════════════════════════════════════════════════════════
s15 = prs.slides.add_slide(blank)
add_rect(s15, 0, 0, W, H, fill=GREY)
add_title_bar(s15, "Testing — 14 Automated + 28 Manual Tests", "TestPortal.java · All tests PASS")

# Left: test categories
cats = [
    ("Positive Testing", GREEN,   "Valid inputs accepted correctly"),
    ("Negative Testing", RED,     "Invalid inputs blocked with clear errors"),
    ("Boundary Testing", ORANGE,  "Edge values (age=18/40, dist=1/100, pace=3/10)"),
    ("Exception Testing", BLUE,   "All exceptions caught, messages readable"),
    ("Role-Based Testing", NAVY,  "Admin sees admin menu, Runner sees runner menu"),
    ("Integration Testing", RGBColor(0x8E,0x24,0xAA), "Register→Verify→Login→Register Race→Log Training"),
    ("Regression Testing", RGBColor(0x00,0x89,0x7B), "All 14 automated tests re-run after every change"),
]
for i, (name, color, desc) in enumerate(cats):
    cy = Inches(1.35 + i * 0.78)
    icon_b = add_rect(s15, Inches(0.25), cy, Inches(0.38), Inches(0.62), fill=color)
    icon_b.line.fill.background()
    add_rect(s15, Inches(0.67), cy, Inches(4.2), Inches(0.62), fill=WHITE).line.color.rgb = color
    add_text(s15, name, Inches(0.72), cy+Inches(0.05), Inches(4.0), Inches(0.3),
             size=11, bold=True, color=color)
    add_text(s15, desc, Inches(0.72), cy+Inches(0.34), Inches(4.0), Inches(0.24),
             size=9, color=NAVY)

# Right: test result table
tests = [
    ("TC-001","Valid Signup","✅ PASS"), ("TC-002","Dup Email","✅ PASS"),
    ("TC-003","Weak Password","✅ PASS"), ("TC-004","Age < 18","✅ PASS"),
    ("TC-005","Unverified Login","✅ PASS"), ("TC-006","5x Lockout","✅ PASS"),
    ("TC-007","Future Training","✅ PASS"), ("TC-008","Fast Pace","✅ PASS"),
    ("TC-009","10x Outlier Flag","✅ PASS"), ("TC-010","<5 Logs Eligibility","✅ PASS"),
    ("TC-011","Double Booking","✅ PASS"), ("TC-012","Late Penalty","✅ PASS"),
    ("TC-013","Spam Detection","✅ PASS"), ("TC-014","URL Routing","✅ PASS"),
]

hx = Inches(5.3)
for j, (hdr, hw) in enumerate(zip(["Test ID","Scenario","Result"],[Inches(1.1),Inches(2.5),Inches(0.85)])):
    hb = add_rect(s15, hx, Inches(1.27), hw, Inches(0.28), fill=NAVY)
    hb.line.fill.background()
    add_text(s15, hdr, hx+Inches(0.04), Inches(1.28), hw-Inches(0.06), Inches(0.24),
             size=9, bold=True, color=WHITE)
    hx += hw

for ri, (tid, scen, res) in enumerate(tests):
    ry = Inches(1.55 + ri * 0.37)
    bg = WHITE if ri % 2 == 0 else LBLUE
    rx = Inches(5.3)
    for cell, hw in zip([tid,scen,res],[Inches(1.1),Inches(2.5),Inches(0.85)]):
        cb = add_rect(s15, rx, ry, hw, Inches(0.36), fill=bg)
        cb.line.color.rgb = RGBColor(0xCC,0xCC,0xCC)
        cb.line.width = Pt(0.4)
        clr = GREEN if "PASS" in cell else NAVY
        add_text(s15, cell, rx+Inches(0.04), ry+Inches(0.05), hw-Inches(0.06), Inches(0.26),
                 size=8.5, bold=("PASS" in cell), color=clr)
        rx += hw

# Summary badge
add_rect(s15, Inches(9.05), Inches(1.27), Inches(4.0), Inches(5.9), fill=WHITE).line.color.rgb = GREEN
add_text(s15, "📊  RESULTS", Inches(9.15), Inches(1.38), Inches(3.8), Inches(0.3),
         size=11, bold=True, color=NAVY)
results = [("Automated Tests","14/14 ✅"),("Manual Tests","28/28 ✅"),
           ("Compilation","0 errors ✅"),("Runtime Crashes","0 ✅"),("Defects Fixed","1 (DA-001) ✅")]
for i, (lbl, val) in enumerate(results):
    ry = Inches(1.75 + i * 0.9)
    add_text(s15, lbl, Inches(9.15), ry, Inches(2.0), Inches(0.3), size=10, color=NAVY)
    add_text(s15, val, Inches(10.5), ry, Inches(2.5), Inches(0.3), size=11, bold=True, color=GREEN)

slide_number(s15, 15)
notes(s15, "TestPortal.java contains 14 automated test cases covering all major business rules. "
      "Run with: javac TestPortal.java && java TestPortal. "
      "All 14 pass including the fixed DA-001 defect where a LaTeX symbol was in an error message. "
      "Manual testing covered boundaries, role separation, and integration flows.")

# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 16 — DEFECT ANALYSIS
# ══════════════════════════════════════════════════════════════════════════════
s16 = prs.slides.add_slide(blank)
add_rect(s16, 0, 0, W, H, fill=GREY)
add_title_bar(s16, "Defect Analysis & Prevention", "How the project guards against common bugs")

defects = [
    ("🚫 Duplicate Email","MEDIUM", GREEN, "Prevented",
     "usedEmails.contains(email) checked before creating account.\n"
     "HashSet guarantees O(1) duplicate detection."),
    ("🔑 Wrong Password","HIGH", GREEN, "Prevented",
     "failedAttempts counter incremented on each failure.\n"
     "5 failures → 30-min lockout via lockedUntil timestamp."),
    ("📅 Future Training Date","MEDIUM", GREEN, "Prevented",
     "date.isAfter(LocalDate.now()) throws Exception.\n"
     "Validated before creating TrainingLog object."),
    ("📋 Duplicate Log","LOW", GREEN, "Prevented",
     "Loop checks same runnerId + date + distance.\n"
     "Throws Exception before saving."),
    ("🔐 Unauthorized Access","HIGH", GREEN, "Prevented",
     "After login, role is checked. RUNNER cannot reach\n"
     "adminMenu() code path — menus are completely separate."),
    ("⚡ Race Validation Bypass","MEDIUM", GREEN, "Prevented",
     "7-day future, unique name, 24h edit, 48h registration\n"
     "all enforced in RaceMgr before saving."),
    ("💥 Input Crash","HIGH", GREEN, "Prevented",
     "readInt() wraps Integer.parseInt() in try-catch.\n"
     "readStr() loops until non-empty input received."),
    ("⚠️ LaTeX in Error Msg (DA-001)","LOW", ORANGE, "Fixed",
     "Error message contained $\\ge 5$ LaTeX notation.\n"
     "Fixed to plain English: 'at least 5'."),
]

cols = 2
for i, (defect, sev, color, status, explanation) in enumerate(defects):
    col = i % cols
    row = i // cols
    dx = Inches(0.25 + col * 6.55)
    dy = Inches(1.3 + row * 1.45)
    dw, dh = Inches(6.3), Inches(1.35)

    header = add_rect(s16, dx, dy, dw, Inches(0.38), fill=color)
    header.line.fill.background()
    add_text(s16, defect, dx+Inches(0.08), dy+Inches(0.06), Inches(4.5), Inches(0.28),
             size=11, bold=True, color=WHITE)
    sev_b = add_rect(s16, dx+dw-Inches(1.2), dy+Inches(0.04), Inches(1.1), Inches(0.28),
                     fill=WHITE if color != ORANGE else ORANGE)
    sev_b.line.fill.background()
    add_text(s16, status, dx+dw-Inches(1.15), dy+Inches(0.06), Inches(1.05), Inches(0.24),
             size=8.5, bold=True, color=color, align=PP_ALIGN.CENTER)

    body = add_rect(s16, dx, dy+Inches(0.38), dw, dh-Inches(0.38), fill=WHITE)
    body.line.color.rgb = color
    body.line.width = Pt(0.75)
    add_text(s16, explanation, dx+Inches(0.1), dy+Inches(0.45), dw-Inches(0.15), dh-Inches(0.5),
             size=9.5, color=NAVY)

slide_number(s16, 16)
notes(s16, "8 potential defect scenarios reviewed. 7 are actively prevented by existing code. "
      "DA-001 (LaTeX notation in error message) was the only real defect found during the review — "
      "it was identified and fixed before this presentation. Zero critical defects remain.")

# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 17 — PROJECT FOLDER STRUCTURE
# ══════════════════════════════════════════════════════════════════════════════
s17 = prs.slides.add_slide(blank)
add_rect(s17, 0, 0, W, H, fill=GREY)
add_title_bar(s17, "Project Folder Structure", "Complete repository layout on GitHub")

folders = [
    ("📁","marathon-runner-portal/   (project root)", NAVY, WHITE, True, 0),
    ("📂","  src/", BLUE, WHITE, True, 1),
    ("📄","    MarathonPortal.java", RGBColor(0x00,0x89,0x7B), WHITE, False, 2),
    ("📄","    TestPortal.java", RGBColor(0x00,0x89,0x7B), WHITE, False, 2),
    ("📂","  exports/", ORANGE, WHITE, True, 1),
    ("📄","    users_YYYYMMDD_HHMMSS.csv", ORANGE, WHITE, False, 2),
    ("📄","    races_YYYYMMDD_HHMMSS.csv", ORANGE, WHITE, False, 2),
    ("📄","    training_*.csv", ORANGE, WHITE, False, 2),
    ("📂","  logs/", RED, WHITE, True, 1),
    ("📄","    security.log", RED, WHITE, False, 2),
    ("📄","  README.md", LTGREY, NAVY, False, 1),
    ("📄","  PROGRAM_STRUCTURE.md", LTGREY, NAVY, False, 1),
    ("📄","  ARCHITECTURE.md", LTGREY, NAVY, False, 1),
    ("📄","  USER_STORY_MAPPING.md", LTGREY, NAVY, False, 1),
    ("📄","  CODE_EXPLANATION.md", LTGREY, NAVY, False, 1),
    ("📄","  TEST_REPORT.md  +  DEFECT_ANALYSIS.md", LTGREY, NAVY, False, 1),
    ("📄","  REVIEW_PREPARATION.md  +  PROJECT_FLOW.md", LTGREY, NAVY, False, 1),
    ("📄","  .gitignore", LTGREY, NAVY, False, 1),
]

for i, (ic, label, bg, fg, is_dir, level) in enumerate(folders):
    y = Inches(1.28 + i * 0.36)
    h = Inches(0.34)
    block = add_rect(s17, Inches(0.25), y, Inches(7.0), h, fill=bg if is_dir else LBLUE)
    block.line.color.rgb = bg
    block.line.width = Pt(0.5)
    lbl_color = fg if is_dir else NAVY
    add_text(s17, f"{ic}  {label}", Inches(0.35 + level*0.2), y+Inches(0.04),
             Inches(6.8), Inches(0.26), size=9.5, bold=is_dir, color=lbl_color)

# Right: descriptions
descs = [
    (Inches(7.5), Inches(1.35), "📂  src/", "Contains the two Java source files.\nThe only files needed to run the project."),
    (Inches(7.5), Inches(2.8), "📂  exports/", "Auto-created when the app starts.\nAll CSV reports are timestamped and saved here."),
    (Inches(7.5), Inches(4.25), "📂  logs/", "Auto-created on startup.\nsecurity.log records all failed login attempts."),
    (Inches(7.5), Inches(5.7), "📄  Documentation", "9 markdown files covering architecture,\nflow, testing, review prep, and more."),
]
for dx, dy, title, body in descs:
    add_rect(s17, dx, dy, Inches(5.5), Inches(1.1), fill=WHITE).line.color.rgb = BLUE
    add_text(s17, title, dx+Inches(0.1), dy+Inches(0.08), Inches(5.3), Inches(0.3),
             size=11, bold=True, color=BLUE)
    add_text(s17, body, dx+Inches(0.1), dy+Inches(0.4), Inches(5.3), Inches(0.65),
             size=10, color=NAVY)

slide_number(s17, 17)
notes(s17, "Repository structure: src/ has the two Java files (main app + test suite). "
      "exports/ and logs/ are auto-created by the application on startup. "
      "Root contains 9 documentation markdown files covering every aspect of the project. "
      ".gitignore excludes .class files, exports/, and logs/ from version control.")

# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 18 — CODE WALKTHROUGH
# ══════════════════════════════════════════════════════════════════════════════
s18 = prs.slides.add_slide(blank)
add_rect(s18, 0, 0, W, H, fill=GREY)
add_title_bar(s18, "Code Walkthrough", "How data flows through the application")

steps = [
    ("1", "Keyboard Input", NAVY,
     "readStr() / readInt() / readDouble()\n"
     "Loop until valid. Catches NumberFormatException.\n"
     "Prevents any crash from bad user input."),
    ("2", "Menu Method Called", BLUE,
     "e.g. addTraining() collects date, distance,\n"
     "duration, notes from user.\n"
     "Wraps manager call in try-catch."),
    ("3", "Manager Validates", GREEN,
     "TrainingMgr.addLog() checks:\n"
     "no future, 1-100km, pace 3-10,\n"
     "no duplicate. Throws if invalid."),
    ("4", "Object Created", ORANGE,
     "new TrainingLog(id, runnerId, date, dist, dur, notes)\n"
     "Pace auto-calculated in constructor.\n"
     "flagged field set if outlier."),
    ("5", "Collection Updated", RGBColor(0x8E,0x24,0xAA),
     "trainingLogs.add(newLog)\n"
     "Static list updated in-place.\n"
     "logIdCount++ increments for next ID."),
    ("6", "Console Output", RGBColor(0x34,0xA8,0x53),
     "ok(\"Training workout saved!\")\n"
     "println(\"Log ID: \" + log.id)\n"
     "warn() if flagged as outlier."),
]

for i, (num, title, color, detail) in enumerate(steps):
    col = i % 3
    row = i // 3
    cx = Inches(0.25 + col * 4.35)
    cy = Inches(1.3 + row * 2.8)
    cw, ch = Inches(4.1), Inches(2.6)

    # Number badge
    badge = add_rect(s18, cx, cy, Inches(0.52), Inches(0.52), fill=color)
    badge.line.fill.background()
    add_text(s18, num, cx+Inches(0.04), cy+Inches(0.04), Inches(0.45), Inches(0.42),
             size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Title
    add_text(s18, title, cx+Inches(0.6), cy+Inches(0.1), cw-Inches(0.65), Inches(0.35),
             size=13, bold=True, color=color)

    # Body
    body = add_rect(s18, cx, cy+Inches(0.58), cw, ch-Inches(0.58), fill=WHITE)
    body.line.color.rgb = color
    body.line.width = Pt(0.75)
    add_text(s18, detail, cx+Inches(0.1), cy+Inches(0.65), cw-Inches(0.15), ch-Inches(0.7),
             size=10, color=NAVY)

    # Arrow between steps
    if col < 2:
        add_text(s18, "→", cx+cw+Inches(0.04), cy+Inches(0.9), Inches(0.3), Inches(0.3),
                 size=18, bold=True, color=color, align=PP_ALIGN.CENTER)

slide_number(s18, 18)
notes(s18, "Six-step data flow: Input → Menu → Manager Validates → Object Created → Collection Updated → Output. "
      "This cycle repeats for every user action. The key insight is that all data mutation happens in Step 5 "
      "and all output happens in Step 6 — clear separation of concerns even within a single method.")

# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 19 — REVIEW PREPARATION
# ══════════════════════════════════════════════════════════════════════════════
s19 = prs.slides.add_slide(blank)
add_rect(s19, 0, 0, W, H, fill=GREY)
add_title_bar(s19, "Review Preparation — Q & A", "Expected invigilator questions and model answers")

qas = [
    ("Why ArrayList?",
     "Ordered, dynamically sized, supports iteration and sort. "
     "Best for sequential data like logs and race lists."),
    ("Why HashMap for users?",
     "O(1) average lookup by email key. Without it, every login "
     "would require O(n) loop through users list."),
    ("Why HashSet for emails?",
     "O(1) duplicate check. Adding email to set after registration "
     "prevents duplicate accounts instantly."),
    ("Why OOP?",
     "Encapsulation keeps data and logic together. Abstraction hides "
     "complexity. Makes code readable and maintainable."),
    ("Why Exception handling?",
     "Prevents crashes. Manager throws descriptive exceptions, "
     "menus catch them and show user-friendly error messages."),
    ("Why no Database?",
     "Project scope is Core Java console only. No JDBC, no MySQL. "
     "In-memory collections reset on exit — expected behaviour."),
    ("Why no Spring/REST?",
     "Spring requires external dependencies and web server. "
     "This project uses only java.util and java.time — pure JDK."),
    ("How is security handled?",
     "5-attempt lockout (30 min) via LocalDateTime. Failed attempts "
     "logged to security.log with timestamp and reason."),
    ("How does role separation work?",
     "After login, user.role is checked. ADMIN calls showAdminMenu(), "
     "RUNNER calls showRunnerMenu(). Paths never cross."),
    ("What if user enters letters for age?",
     "readInt() catches NumberFormatException and re-prompts. "
     "No crash possible — all input is wrapped in while(true)."),
]

col_w = Inches(6.2)
for i, (q, a) in enumerate(qas):
    col = i % 2
    row = i // 2
    qx = Inches(0.25 + col * 6.55)
    qy = Inches(1.3 + row * 1.18)
    add_rect(s19, qx, qy, col_w, Inches(1.1), fill=WHITE).line.color.rgb = BLUE

    q_block = add_rect(s19, qx, qy, col_w, Inches(0.35), fill=NAVY)
    q_block.line.fill.background()
    add_text(s19, f"Q:  {q}", qx+Inches(0.08), qy+Inches(0.05), col_w-Inches(0.12), Inches(0.28),
             size=10, bold=True, color=WHITE)
    add_text(s19, f"A:  {a}", qx+Inches(0.08), qy+Inches(0.4), col_w-Inches(0.12), Inches(0.65),
             size=9.5, color=NAVY)

slide_number(s19, 19)
notes(s19, "10 most likely invigilator questions with model answers. Key points to emphasise: "
      "ArrayList for ordered lists, HashMap for O(1) lookup, HashSet for O(1) duplicate check, "
      "Exception handling for user safety, role-based routing for security, readInt() for crash prevention.")

# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 20 — CONCLUSION
# ══════════════════════════════════════════════════════════════════════════════
s20 = prs.slides.add_slide(blank)
add_rect(s20, 0, 0, W, H, fill=NAVY)
add_rect(s20, 0, Inches(2.8), W, Inches(0.05), fill=BLUE)

add_text(s20, "SPRINT 1 — COMPLETE", Inches(0.5), Inches(0.7), Inches(12.3), Inches(0.7),
         size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s20, "Marathon Runner Portal  ·  Core Java 17  ·  Console Application",
         Inches(0.5), Inches(1.42), Inches(12.3), Inches(0.45),
         size=15, color=ACCENT, align=PP_ALIGN.CENTER)

achievements = [
    ("18/20","User Stories\nImplemented"),
    ("14/14","Automated\nTests PASS"),
    ("1505","Lines of\nClean Code"),
    ("9","Documentation\nFiles"),
    ("0","Critical\nDefects"),
]
for i, (num, lbl) in enumerate(achievements):
    ax = Inches(0.6 + i * 2.42)
    add_rect(s20, ax, Inches(3.1), Inches(2.2), Inches(1.7), fill=BLUE).line.fill.background()
    add_text(s20, num, ax+Inches(0.05), Inches(3.2), Inches(2.1), Inches(0.9),
             size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s20, lbl, ax+Inches(0.05), Inches(4.1), Inches(2.1), Inches(0.55),
             size=10.5, color=ACCENT, align=PP_ALIGN.CENTER)

add_text(s20, "LEARNING OUTCOMES", Inches(0.5), Inches(5.1), Inches(12.2), Inches(0.35),
         size=13, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
outcomes = "OOP Principles  ·  Layered Architecture  ·  Java Collections  ·  Exception Handling  ·  File I/O  ·  Time API  ·  Agile Development"
add_text(s20, outcomes, Inches(0.5), Inches(5.5), Inches(12.2), Inches(0.35),
         size=12, color=WHITE, align=PP_ALIGN.CENTER)

add_text(s20, "FUTURE SCOPE",
         Inches(0.5), Inches(6.0), Inches(12.2), Inches(0.3),
         size=12, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
future = "Top 3 Leaderboard  ·  Forum Posts Export  ·  Force Race Cancellation with Refund  ·  Admin Activity Audit Trail  ·  File-Based Persistence"
add_text(s20, future, Inches(0.5), Inches(6.35), Inches(12.2), Inches(0.3),
         size=10.5, color=RGBColor(0xCC,0xDD,0xFF), align=PP_ALIGN.CENTER)

add_text(s20, "Thank You  🙏",
         Inches(0.5), Inches(6.85), Inches(12.2), Inches(0.45),
         size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

slide_number(s20, 20)
notes(s20, "Closing slide. Key achievements: 18 user stories implemented, 14 automated tests all pass, "
      "1505 lines of well-structured code, 9 documentation files, 0 critical defects. "
      "Learning outcomes: OOP, Collections, Exception Handling, File I/O, Java Time API, Agile. "
      "Future scope: leaderboard, forum export, race cancellation refund, admin audit trail.")

# ── SAVE ───────────────────────────────────────────────────────────────────────
output = r"c:\Users\USER\OneDrive\Desktop\Marathon Runner Portal\Marathon_Runner_Portal_Sprint1.pptx"
prs.save(output)
print(f"✅  Presentation saved → {output}")
print(f"    Slides: 20  |  Format: 16:9  |  Theme: Blue + Navy")
